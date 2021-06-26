from pptx import Presentation
import os, glob

def scanFolder():
    """
    Scans the bin directory, adds all the files in the that directory to a list,
    then filters out all non-image files.
    """
    #common image extensions
    img_Ext = ['.png','.jpg','.jpeg']
    #Changes from current working directory to bin
    os.chdir("bin/")

    #Adds all the files in the bin directory in a list
    files = os.listdir('.')
    imgs = []


    for file in files:
        for ext in img_Ext:
            #filtering out files that do not have an image extension e.g .png or .jpeg
            if ext in file:
                imgs.append(file)
    os.chdir("../")
    return(imgs)


def add_Img_To_File(images):
    """
    Receives a list of image names as an argument, then it loops through each image whilee 
    adding each image to the .pptx file
    """

    pres = Presentation()
    #The layouts in Powerpoint range from 0-8
    layout = pres.slide_layouts[8] #This is the 8th layout

    for img in images:
        slide = pres.slides.add_slide(layout)
        placeHolder = slide.placeHolder[1]
        picture = placeHolder.insert_picture("bin/" + img)

    pres.save("testPres.pptx")

if __name__=="__main__":
    add_Img_To_File(scanFolder())
    print("Done!")
